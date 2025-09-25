# app/laudo_pptx.py
# =============================================================================
# Geração de Laudo (PPTX + PDF) com gráficos e placeholders
# - Renderiza gráficos como PNGs transparentes (matplotlib)
# - Preenche textos e imagens no template PPTX (python-pptx)
# - Converte para PDF via LibreOffice Impress (com suporte a fontes do repo)
# - Fallback: se falhar o PDF vetorial, exporta PNGs dos slides e monta um PDF
# =============================================================================

import io, os, re, math, json, shutil, subprocess, tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Sequence, Union, Tuple

# gráficos
import matplotlib.pyplot as plt

# pptx / imagens
from pptx import Presentation
from pptx.dml.color import RGBColor, MSO_COLOR_TYPE
from pptx.util import Inches
from PIL import Image

import uuid
# conversão PPTX -> PDF (utils)
from utils.pptx_pdf import convert_pptx_to_pdf

# Prefixo de pasta no /tmp (ex.: /tmp/laudo_<UUID>)
LAUDO_PREFIX = "laudo_"

# Regex para encontrar {{chave}} (com ou sem espaços)
REX = re.compile(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}")

# Pastas base para localizar templates
BASE_DIR = Path(__file__).resolve().parent.parent
TPL_DIR = BASE_DIR / "templates"

# -----------------------------------------------------------------------------
# Helpers numéricos e rótulos
# -----------------------------------------------------------------------------
PT_BR_SHORT = ["Jan","Fev","Mar","Abr","Mai","Jun","Jul","Ago","Set","Out","Nov","Dez"]

def _periodo_to_label(periodo: str) -> str:
    """
    Converte "AAAA-Q" para "Qº Tri\\nAAAA". Ex.: "2025-1" -> "1º Tri\\n2025".
    """
    s = str(periodo).strip()
    if "-" not in s:
        return s
    ano, tri = s.split("-", 1)
    try:
        q = int(tri)
    except ValueError:
        return s
    return f"{q}º Tri\n{ano}"

def _to_float(v) -> float:
    """
    Converte string/num em float (aceita vírgula como decimal).
    Retorna NaN em caso de falha.
    """
    if v is None or str(v).strip() == "":
        return float("nan")
    try:
        return float(str(v).replace(",", "."))
    except ValueError:
        return float("nan")

def _coerce_float(x) -> float:
    """
    Conversor mais "esperto" para BRL/strings variadas: tenta normalizar
    12.345,67 -> 12345.67 ; 12,34 -> 12.34 ; "1234" -> 1234.0
    """
    if x is None:
        return float("nan")
    if isinstance(x, (int, float)):
        return float(x)
    s = str(x).strip()
    if "," in s:
        s = s.replace(".", "").replace(",", ".")
    else:
        parts = s.split(".")
        if len(parts) > 2:
            s = "".join(parts[:-1]) + "." + parts[-1]
        elif len(parts) == 2 and len(parts[1]) <= 2:
            s = ".".join(parts)
        else:
            s = s.replace(".", "")
    try:
        return float(s)
    except Exception:
        return float("nan")

def _meses_labels_2linhas(inicio_ym: str, n: int = 12):
    """
    Gera rótulos 'Mês\\nAno' a partir de um 'AAAA-MM' inicial, por n pontos.
    """
    ano, mes = inicio_ym.split("-")
    a = int(ano); m = int(mes)
    return [f"{PT_BR_SHORT[(m-1+i)%12]}\n{a + (m-1+i)//12}" for i in range(n)]

# -----------------------------------------------------------------------------
# Gráficos (PNG transparente, sem borda de figura/eixos)
# -----------------------------------------------------------------------------
def grafico_trimestral_png(
    data: List[Dict[str, Any]],
    label_key: str = "periodo",
    anunciados_key: str = "anuncios",  # laranja
    vendidos_key: str = "vendidos",    # roxo
    ylim: Optional[Sequence[float]] = None,
    dpi: int = 200,
    figsize: Tuple[float, float] = (9, 4),
) -> bytes:
    """
    Gráfico de linhas (anunciados vs vendidos) com fundo transparente.
    Retorna bytes PNG (bbox tight) para caber limpo no PPTX.
    """
    # Eixo X
    labels_raw = [str(r.get(label_key, "")) for r in data]
    x_labels = [_periodo_to_label(s) for s in labels_raw]
    X = list(range(len(x_labels)))

    # Séries
    y_anun = [_to_float(r.get(anunciados_key)) for r in data]  # laranja
    y_vend = [_to_float(r.get(vendidos_key)) for r in data]    # roxo

    # Escala Y (múltiplos de 20, mínimo 80) se não vier definida
    if ylim is None:
        vals = [v for v in (y_anun + y_vend) if not (isinstance(v, float) and math.isnan(v))]
        if not vals:
            raise ValueError("Sem valores numéricos.")
        ymax = max(80, math.ceil(max(vals) / 20) * 20)
        y_min, y_max = 0, ymax
    else:
        y_min, y_max = ylim

    # Figura/Axes sem fundo e sem spines
    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
    fig.patch.set_alpha(0)         # sem fundo da figura
    ax.set_facecolor("none")       # sem fundo do eixo
    ax.grid(True, which="major", axis="both", linestyle="-", linewidth=1, color="#D9D9D9", alpha=0.35)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.tick_params(axis="x", colors="#46484C", labelsize=9)
    ax.tick_params(axis="y", colors="#46484C", labelsize=9)

    ax.set_ylim(y_min, y_max)
    ax.set_yticks(list(range(int(y_min), int(y_max) + 1, 20)))
    ax.set_xticks(X)
    ax.set_xticklabels(x_labels)

    # Linhas (ordem de desenho)
    ax.plot(X, y_vend, marker="o", markersize=6, linewidth=2.5, color="#7a2be2")  # roxo
    ax.plot(X, y_anun, marker="o", markersize=6, linewidth=2.5, color="#f39c12")  # laranja

    plt.tight_layout(pad=0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0, transparent=True)
    plt.close(fig)
    return buf.getvalue()

def grafico_area_mensal_png(
    valores: List[Union[int, float, str]],
    inicio_ym: str = "2023-08",
    moeda_prefix: str = "R$ ",
    dpi: int = 200,
    figsize: Tuple[float, float] = (12, 3.2),
) -> bytes:
    """
    Gráfico de área mensal (12 pontos), fundo transparente.
    """
    if len(valores) != 12:
        raise ValueError("Forneça 12 valores (um por mês).")

    ys = [_coerce_float(v) for v in valores]
    xs = list(range(12))
    labels = _meses_labels_2linhas(inicio_ym, 12)

    valid = [v for v in ys if not math.isnan(v)]
    if not valid:
        raise ValueError("Sem valores numéricos válidos.")
    step = 5000.0
    ymax = max(step, math.ceil(max(valid) / step) * step)

    purple = "#720d83"
    grid_c  = "#46484C"
    axis_c  = "#e6e6e6"
    txt_c   = "#46484C"

    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")

    # sem vão nas bordas
    ax.set_xlim(0, 11)
    ax.margins(x=0)

    # grade / eixos
    ax.grid(True, which="major", axis="both", linestyle="-", linewidth=1, color=grid_c, alpha=0.25)
    for s in ax.spines.values():
        s.set_color(axis_c); s.set_linewidth(1)
    ax.tick_params(axis="x", colors=txt_c, labelsize=10)
    ax.tick_params(axis="y", colors=txt_c, labelsize=10)

    # Y formatado em BRL (ou prefixo fornecido)
    ax.set_ylim(0, ymax)
    yticks = list(range(0, int(ymax) + 1, int(step)))
    ax.set_yticks(yticks)
    ax.set_yticklabels([f"{moeda_prefix}{int(v):,}".replace(",", ".") for v in yticks])

    # X rótulos
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)

    # área + linha + pontos
    ax.fill_between(xs, ys, 0, color=purple, alpha=0.85, linewidth=0)
    ax.plot(xs, ys, color=purple, linewidth=2.5, marker="o", markersize=5)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", transparent=True)
    plt.close(fig)
    return buf.getvalue()

# -----------------------------------------------------------------------------
# Helpers de texto/fonte para PPTX
# -----------------------------------------------------------------------------
def _snapshot_font(run_font):
    """Captura nome/tamanho/estilos e cor (RGB ou SchemeColor + brilho)."""
    name = run_font.name
    size = run_font.size
    bold = run_font.bold
    italic = run_font.italic
    color_info = None
    cf = run_font.color
    if cf is not None and cf.type is not None:
        try:
            if cf.type == MSO_COLOR_TYPE.RGB and cf.rgb:
                color_info = ("rgb", cf.rgb, None)
            elif cf.type == MSO_COLOR_TYPE.SCHEME and cf.theme_color is not None:
                br = None
                try:
                    br = cf.brightness
                except Exception:
                    pass
                color_info = ("scheme", cf.theme_color, br)
        except Exception:
            pass
    return (name, size, bold, italic, color_info)

def _restore_font(run_font, snap):
    """Restaura nome/tamanho/estilos e cor (RGB ou SchemeColor + brilho)."""
    name, size, bold, italic, color_info = snap
    if name: run_font.name = name
    if size: run_font.size = size
    if bold is not None: run_font.bold = bold
    if italic is not None: run_font.italic = italic
    if color_info:
        kind, val, br = color_info
        try:
            if kind == "rgb" and val:
                run_font.color.rgb = RGBColor(val[0], val[1], val[2])
            elif kind == "scheme" and val is not None:
                run_font.color.theme_color = val
                if br is not None:
                    run_font.color.brightness = br
        except Exception:
            pass

def _replace_text_in_textframe(text_frame, get_value_fn, img_keys_set):
    """
    Substitui placeholders de TEXTO dentro de um text_frame, preservando
    as propriedades de fonte do run original.
    """
    for p in text_frame.paragraphs:
        for r in p.runs:
            txt = r.text
            if "{{" not in txt:
                continue
            # se o run for exatamente {{chave}} e for de imagem, pula aqui
            full = REX.fullmatch(txt.strip())
            if full and full.group(1) in img_keys_set:
                continue
            snap = _snapshot_font(r.font)
            def repl(m):
                k = m.group(1)
                v = get_value_fn(k)
                return v if v is not None else m.group(0)
            new_txt = REX.sub(repl, txt)
            if new_txt != txt:
                r.text = new_txt
                _restore_font(r.font, snap)

def _find_img_keys(shape, img_keys_set):
    """
    Detecta se a caixa de texto contém APENAS um {{chave}} que
    corresponde a uma imagem. Retorna lista de chaves.
    """
    keys = []
    if hasattr(shape, "text_frame") and shape.text_frame:
        text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs).strip()
        m = REX.fullmatch(text)
        if m:
            k = m.group(1)
            if k in img_keys_set:
                keys.append(k)
    return keys

def _add_image(slide, image_path, box, target_size=None):
    """
    Insere imagem.
    - Se target_size=(w,h) for passado, usa tamanho fixo (EMUs) e centraliza no 'box'.
    - Caso contrário, faz auto-fit "contain" para caber no 'box' mantendo proporção.
    """
    image_path = str(image_path)
    if not Path(image_path).exists():
        raise FileNotFoundError(f"Imagem não encontrada: {image_path}")

    l, t, w, h = box
    if target_size:
        tw, th = target_size
        left = int(l + (w - tw) / 2)
        top  = int(t + (h - th) / 2)
        slide.shapes.add_picture(image_path, left, top, width=int(tw), height=int(th))
        return

    # auto-fit (contain)
    img_w, img_h = Image.open(image_path).size  # px
    scale = min(w / img_w, h / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)
    left = int(l + (w - new_w) / 2)
    top  = int(t + (h - new_h) / 2)
    slide.shapes.add_picture(image_path, left, top, width=new_w, height=new_h)

def _walk(slide, shape, images_to_place, get_value_fn, img_keys_set):
    """
    Percorre shapes (inclusive grupos e tabelas) para substituir textos e
    coletar caixas que são placeholders de imagem.
    """
    # Tabela
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    _replace_text_in_textframe(cell.text_frame, get_value_fn, img_keys_set)

    # Caixa de texto
    if hasattr(shape, "text_frame") and shape.text_frame:
        img_keys = _find_img_keys(shape, img_keys_set)
        if img_keys:
            images_to_place.append((slide, shape, img_keys))
        else:
            _replace_text_in_textframe(shape.text_frame, get_value_fn, img_keys_set)

    # Grupo
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            _walk(slide, s, images_to_place, get_value_fn, img_keys_set)

# -----------------------------------------------------------------------------
# Conversão para PDF via LibreOffice (com fontes do repo) + Fallback PNG->PDF
# -----------------------------------------------------------------------------
def _candidate_fonts_dirs() -> list[Path]:
    """
    Procura diretórios de fontes dentro do repositório para expor via FONTCONFIG_FILE.
    Ajuste aqui se você usar outras famílias além de Nunito.
    """
    here = Path(__file__).resolve().parent
    repo_root = here.parent
    candidates = [
        repo_root / "fonts" / "nunito" / "static",
        repo_root / "fonts" / "nunito",
        repo_root / "fonts",
    ]
    return [p for p in candidates if p.exists() and any(p.glob("*.ttf"))]

def _write_fonts_conf(out_dir: Path, font_dirs: list[Path]) -> Path:
    """
    Gera um fonts.conf temporário apontando para os diretórios informados.
    Retorna o caminho do arquivo gerado.
    """
    fc_dir = out_dir / ".fontconfig"
    fc_dir.mkdir(parents=True, exist_ok=True)
    conf_path = fc_dir / "fonts.conf"
    # precisa ser caminho ABSOLUTO
    dirs_xml = "\n".join(f"  <dir>{str(d.resolve())}</dir>" for d in font_dirs)
    conf_xml = f"""<?xml version="1.0"?>
<!DOCTYPE fontconfig SYSTEM "fonts.dtd">
<fontconfig>
{dirs_xml}
</fontconfig>"""
    conf_path.write_text(conf_xml, encoding="utf-8")
    return conf_path

def _libreoffice_pdf(pptx_path: Path, out_dir: Path):
    """
    PPTX -> PDF via LibreOffice Impress com fontes do repo (Nunito) e profile isolado.
    Gera logs (lo_stdout.log / lo_stderr.log) na pasta do laudo.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    log_out = out_dir / "lo_stdout.log"
    log_err = out_dir / "lo_stderr.log"

    cmd = shutil.which("libreoffice") or shutil.which("soffice") or shutil.which("loffice")
    if not cmd:
        raise RuntimeError("LibreOffice (libreoffice/soffice/loffice) não encontrado no PATH.")

    # Profile dedicado (evita HOME sem permissão)
    lo_profile = out_dir / ".lo_profile"
    lo_profile.mkdir(parents=True, exist_ok=True)
    lo_profile_uri = f"file://{lo_profile.resolve()}"

    # Ambiente estável/headless
    env = os.environ.copy()
    env.setdefault("SAL_USE_VCLPLUGIN", "gen")
    env.setdefault("HOME", "/tmp")
    env.setdefault("XDG_CACHE_HOME", "/tmp/.cache")

    # Fontes do repo
    font_dirs = _candidate_fonts_dirs()
    if font_dirs:
        conf_path = _write_fonts_conf(out_dir, font_dirs)
        env["FONTCONFIG_FILE"] = str(conf_path)
        try:
            subprocess.run(["fc-cache", "-f", "-v"], check=True, env=env,
                           stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except Exception:
            pass

    args = [
        cmd, "--headless", "--nologo", "--nolockcheck", "--norestore", "--nodefault",
        f"-env:UserInstallation={lo_profile_uri}",
        "--convert-to", "pdf:impress_pdf_Export",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    proc = subprocess.run(args, env=env, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    log_out.write_text(proc.stdout or "", encoding="utf-8")
    log_err.write_text(proc.stderr or "", encoding="utf-8")

    pdf_out = out_dir / (pptx_path.with_suffix(".pdf").name)
    if proc.returncode != 0 or not pdf_out.exists():
        raise RuntimeError(
            "Falha na conversão LibreOffice (PPTX).\n"
            f"cmd: {' '.join(args)}\n"
            f"stdout:\n{proc.stdout}\n\nstderr:\n{proc.stderr}\n"
            f"FONTCONFIG_FILE={env.get('FONTCONFIG_FILE')}\n"
            f"UserInstallation={lo_profile_uri}"
        )

def _libreoffice_pngs(pptx_path: Path, out_dir: Path) -> list[Path]:
    """
    Exporta slides como PNG usando o LibreOffice Impress.
    Retorna a lista de PNGs gerados (ordenada). Gera logs separados.
    """
    cmd = shutil.which("libreoffice") or shutil.which("soffice") or shutil.which("loffice")
    if not cmd:
        raise RuntimeError("LibreOffice não encontrado no PATH.")

    env = os.environ.copy()
    env.setdefault("SAL_USE_VCLPLUGIN", "gen")
    env.setdefault("HOME", "/tmp")
    env.setdefault("XDG_CACHE_HOME", "/tmp/.cache")

    # usa o mesmo profile da conversão principal
    lo_profile = out_dir / ".lo_profile"
    lo_profile.mkdir(parents=True, exist_ok=True)
    lo_profile_uri = f"file://{lo_profile.resolve()}"

    if (out_dir / ".fontconfig" / "fonts.conf").exists():
        env["FONTCONFIG_FILE"] = str((out_dir / ".fontconfig" / "fonts.conf").resolve())

    args = [
        cmd, "--headless", "--nologo", "--nolockcheck", "--norestore", "--nodefault",
        f"-env:UserInstallation={lo_profile_uri}",
        "--convert-to", "png:impress_png_Export",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    proc = subprocess.run(args, env=env, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    (out_dir / "lo_png_stdout.log").write_text(proc.stdout or "", encoding="utf-8")
    (out_dir / "lo_png_stderr.log").write_text(proc.stderr or "", encoding="utf-8")

    # LibreOffice costuma nomear como "<entrada>_1.png", "_2.png", ...
    stem = pptx_path.stem
    pngs = sorted(out_dir.glob(f"{stem}*.png"), key=lambda p: p.name.lower())
    if not pngs:
        raise RuntimeError(
            "Nenhum PNG gerado pelo LibreOffice."
            f"\nstdout:\n{proc.stdout}\n\nstderr:\n{proc.stderr}"
        )
    return pngs

def _pngs_to_pdf(png_paths: list[Path], out_pdf: Path):
    """
    Junta PNGs em um único PDF (perde seleção de texto, mas é pixel-perfect).
    """
    images = [Image.open(str(p)).convert("RGB") for p in png_paths]
    first, rest = images[0], images[1:]
    first.save(str(out_pdf), save_all=True, append_images=rest)

# -----------------------------------------------------------------------------
# Núcleo: monta PPTX a partir do template, injeta gráficos/imagens, salva e gera PDF
# -----------------------------------------------------------------------------
def render_laudo(
    template_path: Union[str, Path],
    text_vars: Dict[str, str],
    aliases: Optional[Dict[str, str]] = None,
    chart1: Optional[Dict[str, Any]] = None,
    chart2: Optional[Dict[str, Any]] = None,
    images_bindings: Optional[Dict[str, Union[str, Tuple[str, float, float]]]] = None,
    chart_slots: Optional[Dict[str, str]] = None,
    out_basename: Optional[str] = None,
    force_font_family: Optional[str] = "Nunito",  # força família ao final (opcional)
) -> Dict[str, str]:
    """
    Gera PPTX (e PDF) a partir de um template e payload.
    - template_path: caminho do template .pptx (ou usa fallback em templates/laudo-imogo.pptx)
    - text_vars: dict com variáveis de texto -> substituem {{chave}}
    - aliases: map de alias -> chave real (quando template usa outro nome)
    - chart1/2: configurações dos gráficos (ver chamadas acima)
    - images_bindings: dict com bindings de imagens ("foto": ["path", w_in, h_in] ou "foto": "path")
    - chart_slots: mapeia "chart1"/"chart2" para nomes de placeholders no PPTX (ex.: "grafico_01")
    - out_basename: define o UUID/nome base (senão gera um)
    - force_font_family: se quiser forçar "Nunito" no texto (ou None para não alterar)
    """
    template_path = Path(template_path)
    if not template_path.exists():
        # fallback para templates/laudo-imogo.pptx
        candidate = (TPL_DIR / "laudo-imogo.pptx")
        if candidate.exists():
            template_path = candidate
        else:
            raise FileNotFoundError(f"Template não encontrado: {template_path}")

    rid = out_basename or str(uuid.uuid4())

    # pasta de trabalho em /tmp com prefixo (ex.: /tmp/laudo_<UUID>)
    work = Path(tempfile.gettempdir()) / f"{LAUDO_PREFIX}{rid}"
    work.mkdir(exist_ok=True)

    # 1) Gera gráficos (salva PNGs na pasta de trabalho)
    gen_imgs: Dict[str, Path] = {}
    if chart1:
        png = grafico_trimestral_png(
            data=chart1["data"],
            label_key=chart1.get("label_key", "periodo"),
            anunciados_key=chart1.get("anunciados_key", "anuncios"),
            vendidos_key=chart1.get("vendidos_key", "vendidos"),
            ylim=chart1.get("ylim"),
        )
        p = work / "grafico_trimestres_transparente.png"
        p.write_bytes(png)
        gen_imgs["chart1"] = p

    if chart2:
        png = grafico_area_mensal_png(
            valores=chart2["valores"],
            inicio_ym=chart2.get("inicio_ym", "2023-08"),
            moeda_prefix=chart2.get("moeda_prefix", "R$ "),
        )
        p = work / "grafico2_transparente.png"
        p.write_bytes(png)
        gen_imgs["chart2"] = p

    # 2) Vars e imagens
    VARS = dict(text_vars or {})
    ALIASES = dict(aliases or {})
    IMG_VARS: Dict[str, Union[str, Tuple[Path, float, float]]] = {}
    images_bindings = images_bindings or {}

    # Bindings de imagens explícitos (path + largura/altura em polegadas)
    for k, v in images_bindings.items():
        if isinstance(v, (list, tuple)) and len(v) == 3:
            IMG_VARS[k] = (Path(v[0]), float(Inches(v[1])), float(Inches(v[2])))
        elif isinstance(v, str):
            IMG_VARS[k] = v
        else:
            raise ValueError(f"Imagem inválida para '{k}'")

    # Slots padrão para charts (pode ser sobrescrito pelo payload)
    slots = chart_slots or {"chart1": "grafico_01", "chart2": "grafico_02"}
    if "chart1" in gen_imgs:
        varname = slots.get("chart1", "grafico_01")
        # tamanho padrão no slide (ajuste conforme layout)
        IMG_VARS[varname] = (gen_imgs["chart1"], float(Inches(4.0)), float(Inches(1.8)))
    if "chart2" in gen_imgs:
        varname = slots.get("chart2", "grafico_02")
        IMG_VARS[varname] = (gen_imgs["chart2"], float(Inches(6.5)), float(Inches(1.5)))

    def _get_value(k: str):
        if k in VARS:
            return VARS[k]
        if k in ALIASES and ALIASES[k] in VARS:
            return VARS[ALIASES[k]]
        return None  # mantém placeholder se não tiver valor

    img_keys = set(IMG_VARS.keys())

    # 3) Carrega/varre PPTX
    prs = Presentation(str(template_path))
    to_place = []  # (slide, shape, [keys])
    for slide in prs.slides:
        for shape in slide.shapes:
            _walk(slide, shape, to_place, _get_value, img_keys)

    # 4) Insere imagens nas caixas {{chave}}
    for slide, shape, keys in to_place:
        key = keys[0]   # 1 por caixa
        info = IMG_VARS.get(key)

        # limpa texto da caixa do placeholder
        if hasattr(shape, "text_frame") and shape.text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""

        if info is None:
            continue

        if isinstance(info, tuple):
            img_path, w_emus, h_emus = info
            target = (int(w_emus), int(h_emus))
        else:
            img_path = info
            target = None

        box = (shape.left, shape.top, shape.width, shape.height)
        _add_image(slide, img_path, box, target)

    # (Opcional) Força família de fonte para todo o texto (melhor fidelidade no LO)
    if force_font_family:
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.name = force_font_family
                    if hasattr(shape, "shapes"):
                        for s in shape.shapes:
                            if hasattr(s, "text_frame") and s.text_frame:
                                for p in s.text_frame.paragraphs:
                                    for r in p.runs:
                                        r.font.name = force_font_family
        except Exception:
            # Melhor não falhar por causa disso
            pass

    # Salva PPTX final
    out_pptx = work / f"{rid}.pptx"
    prs.save(out_pptx)

    # 5) Gera PDF (vetorial, se possível), senão fallback PNG->PDF
    pdf_path = work / f"{rid}.pdf"
    pdf_ok = convert_pptx_to_pdf(out_pptx, out_dir=work)

    return {
        "id": rid,
        "dir": str(work),
        "pptx_path": str(out_pptx),
        "pdf_path": str(pdf_path) if pdf_ok and pdf_path.exists() else ""
    }